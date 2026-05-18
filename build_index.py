#!/usr/bin/env python3
"""
扫描 03资料库 下所有常见格式文档 → 抽文本到 .extracted-text/ → 写 资料库索引.json。
支持: MD / PDF / TXT / DOC / DOCX / XLSX / XLS / PPTX / PPT

用法:
  python3 build_index.py            # 增量（已抽过的跳过）
  python3 build_index.py --force    # 强制全部重抽
"""
import json
import os
import re
import subprocess
import sys
import tempfile
from collections import Counter
from pathlib import Path

ROOT = Path("/Users/sunyiting/Library/Mobile Documents/iCloud~md~obsidian/Documents/苍茫云海间/50-项目/主网自动化竞赛/03资料库")
INDEX = ROOT / "资料库索引.json"
TEXT_DIR = ROOT / ".extracted-text"
OLD_TEXT_DIR = ROOT / ".pdf-text"  # 旧目录兼容

SUPPORTED = {".md", ".pdf", ".txt", ".doc", ".docx", ".xlsx", ".xls", ".pptx", ".ppt"}
STOP = set("的 了 和 是 在 有 与 也 或 就 都 等 及 对 中 上 下 个 一 二 三 四 五 之 其 此 这 那 但 而 如 the and or of in to a an for is on with by".split())

FORCE = "--force" in sys.argv

# -------- 提取器 --------

def extract_pdf(p, out):
    subprocess.run(["pdftotext", "-layout", "-enc", "UTF-8", str(p), str(out)],
                   check=True, capture_output=True, timeout=180)

def extract_doc_docx(p, out):
    # 优先 textutil（macOS 自带，docx/doc/rtf 都能转）
    try:
        subprocess.run(["textutil", "-convert", "txt", "-encoding", "UTF-8",
                        "-output", str(out), str(p)],
                       check=True, capture_output=True, timeout=120)
    except Exception:
        # 兜底 antiword（仅 .doc）
        if p.suffix.lower() == ".doc":
            r = subprocess.run(["antiword", str(p)], capture_output=True, timeout=60)
            out.write_bytes(r.stdout)
        else:
            raise

def extract_xlsx(p, out):
    import openpyxl
    wb = openpyxl.load_workbook(p, read_only=True, data_only=True)
    lines = []
    for ws in wb.worksheets:
        lines.append(f"### Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                lines.append("\t".join(cells))
    out.write_text("\n".join(lines), encoding="utf-8")

def extract_xls(p, out):
    import xlrd
    wb = xlrd.open_workbook(p)
    lines = []
    for ws in wb.sheets():
        lines.append(f"### Sheet: {ws.name}")
        for r in range(ws.nrows):
            row = ws.row_values(r)
            cells = [str(c) for c in row if c not in (None, "")]
            if cells:
                lines.append("\t".join(cells))
    out.write_text("\n".join(lines), encoding="utf-8")

def extract_pptx(p, out):
    from pptx import Presentation
    prs = Presentation(p)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"### Slide {i}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text)
    out.write_text("\n".join(lines), encoding="utf-8")

def extract_ppt(p, out):
    # textutil 不直接支持 .ppt；改用 LibreOffice 兜底（如有），否则放弃
    soffice = subprocess.run(["which", "soffice"], capture_output=True, text=True)
    if soffice.returncode != 0:
        raise RuntimeError("需要 LibreOffice 才能转 .ppt；可手动转 .pptx")
    with tempfile.TemporaryDirectory() as td:
        subprocess.run(["soffice", "--headless", "--convert-to", "txt",
                        "--outdir", td, str(p)], check=True, capture_output=True, timeout=180)
        txt = next(Path(td).glob("*.txt"))
        out.write_text(txt.read_text(encoding="utf-8", errors="ignore"), encoding="utf-8")

def extract_txt_md(p, out):
    out.write_text(p.read_text(encoding="utf-8", errors="ignore"), encoding="utf-8")

EXTRACTORS = {
    ".md": extract_txt_md,
    ".txt": extract_txt_md,
    ".pdf": extract_pdf,
    ".docx": extract_doc_docx,
    ".doc": extract_doc_docx,
    ".xlsx": extract_xlsx,
    ".xls": extract_xls,
    ".pptx": extract_pptx,
    ".ppt": extract_ppt,
}

# -------- 元数据提取 --------

def get_pdf_pages(p):
    try:
        r = subprocess.run(["pdfinfo", str(p)], capture_output=True, text=True, timeout=30)
        m = re.search(r"Pages:\s+(\d+)", r.stdout)
        return int(m.group(1)) if m else 0
    except Exception:
        return 0

def parse_md_frontmatter(text):
    """从 MD 文本里抽 frontmatter 的 tags / title。"""
    fm_tags, title = [], None
    m = re.match(r"^---\n(.*?)\n---\n", text, re.DOTALL)
    if m:
        fm = m.group(1)
        t = re.search(r"^tags:\s*\[(.*?)\]", fm, re.MULTILINE)
        if t:
            fm_tags = [x.strip().strip("\"'") for x in t.group(1).split(",") if x.strip()]
        body = text[m.end():]
    else:
        body = text
    h = re.search(r"^#\s+(.+)$", body, re.MULTILINE)
    if h:
        title = h.group(1).strip()
    return fm_tags, title

def make_desc(text, n=500):
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    lines = [l for l in lines if not re.fullmatch(r"\d+", l) and len(l) > 1]
    return "\n".join(lines)[:n].replace("\n", " / ")

def extract_keywords(text, top=15):
    cn = re.findall(r"[一-龥]{2,4}", text)
    en = re.findall(r"[A-Za-z][A-Za-z0-9_]{2,}", text)
    tokens = [t.lower() for t in en if t.lower() not in STOP] + [t for t in cn if t not in STOP]
    return [w for w, _ in Counter(tokens).most_common(top * 2) if len(w) >= 2][:top]

# -------- 主流程 --------

def process(p):
    rel = p.relative_to(ROOT)
    rel_str = str(rel)
    dir_name = rel.parts[0] if len(rel.parts) > 1 else ""
    ext = p.suffix.lower()

    # MD 直接读，不写 sidecar
    if ext == ".md":
        text = p.read_text(encoding="utf-8", errors="ignore")
        tags, title = parse_md_frontmatter(text)
        # 去掉 frontmatter 后再抽 desc
        body = re.sub(r"^---\n.*?\n---\n", "", text, count=1, flags=re.DOTALL)
        return {
            "file": rel_str, "title": title or p.stem, "dir": dir_name,
            "type": "md", "tags": tags, "subject": dir_name,
            "desc": make_desc(body), "keywords": extract_keywords(body),
        }

    # 其余格式 → 抽到 .extracted-text/
    txt_rel = Path(".extracted-text") / rel.with_suffix(".txt")
    txt_path = ROOT / txt_rel
    if FORCE or not txt_path.exists():
        txt_path.parent.mkdir(parents=True, exist_ok=True)
        try:
            EXTRACTORS[ext](p, txt_path)
        except Exception as e:
            print(f"  [FAIL] {rel_str}: {e}", file=sys.stderr)
            return None

    try:
        text = txt_path.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        text = ""

    entry = {
        "file": rel_str, "title": p.stem, "dir": dir_name,
        "type": ext.lstrip("."), "text_file": str(txt_rel),
        "tags": [], "subject": dir_name,
        "desc": make_desc(text) if text else "(抽取失败)",
        "keywords": extract_keywords(text) if text else [],
    }
    if ext == ".pdf":
        entry["pages"] = get_pdf_pages(p)
    return entry

def main():
    # 迁移旧 .pdf-text/ → .extracted-text/（如存在）
    if OLD_TEXT_DIR.exists() and not TEXT_DIR.exists():
        OLD_TEXT_DIR.rename(TEXT_DIR)
        print(f"已迁移 {OLD_TEXT_DIR.name} → {TEXT_DIR.name}")

    files = sorted([p for p in ROOT.rglob("*")
                    if p.is_file() and p.suffix.lower() in SUPPORTED
                    and not p.name.endswith(".excalidraw.md")
                    and not any(part.startswith(".") for part in p.relative_to(ROOT).parts)])
    print(f"待处理文件: {len(files)}")

    entries = []
    counts = Counter()
    for i, p in enumerate(files, 1):
        print(f"[{i}/{len(files)}] {p.relative_to(ROOT)}", flush=True)
        e = process(p)
        if e:
            entries.append(e)
            counts[e["type"]] += 1

    with INDEX.open("w", encoding="utf-8") as f:
        json.dump(entries, f, ensure_ascii=False, indent=2)

    print(f"\n完成: 总 {len(entries)} 条")
    for t, n in counts.most_common():
        print(f"  {t}: {n}")
    print(f"索引: {INDEX}")
    print(f"抽取目录: {TEXT_DIR}")

if __name__ == "__main__":
    main()
